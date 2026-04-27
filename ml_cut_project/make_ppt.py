#!/usr/bin/env python3
"""Generate a clean, black-and-white project summary PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY  = RGBColor(90, 90, 90)
LGRAY = RGBColor(200, 200, 200)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ── helpers ───────────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    return box.text_frame

def para(tf, text, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
         space_after=Pt(6), first=False):
    if first and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    p.space_after = space_after
    p.line_spacing = 1.15

def bullet(tf, text, sz=17, bold=False, level=0, color=BLACK):
    p = tf.add_paragraph()
    # Add manual bullet and indent
    if level == 0:
        p.text = "•  " + text
        p.margin_left = Inches(0.3)
        p.level = 0
    else:
        p.text = "◦  " + text
        p.margin_left = Inches(0.3 + (0.4 * level))
        p.level = level
        
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.space_after = Pt(14)
    p.line_spacing = 1.15

def heading(slide, text, top=Inches(0.35)):
    tf = tb(slide, Inches(0.8), top, Inches(11.7), Inches(0.75))
    para(tf, text, sz=30, bold=True, first=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
          Inches(0.8), top + Inches(0.65), Inches(11.7), Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

def content(slide):
    return tb(slide, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.8))

def deco_bars(slide):
    for y in [0, H - Inches(0.12)]:
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Inches(0.12))
        b.fill.solid(); b.fill.fore_color.rgb = BLACK; b.line.fill.background()

def add_table(slide, data, left, top, width, row_h, col_widths, header=True):
    """data = list of lists; first row is header."""
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(row_h * rows))
    tbl = tbl_shape.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = "Arial"
                p.font.bold = (ri == 0)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
            else:
                cell.fill.fore_color.rgb = WHITE
    return tbl

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); deco_bars(s)
tf = tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.8))
para(tf, "ML-Guided Cut Selection for", sz=38, bold=True, align=PP_ALIGN.CENTER, space_after=Pt(2), first=True)
para(tf, "ASIC Technology Mapping in ABC", sz=38, bold=True, align=PP_ALIGN.CENTER, space_after=Pt(24))

tf2 = tb(s, Inches(1), Inches(3.8), Inches(11.3), Inches(2.5))
para(tf2, "Submitted By", sz=20, bold=True, align=PP_ALIGN.CENTER, color=GRAY, first=True, space_after=Pt(14))
para(tf2, "Hemang Gautam", sz=26, align=PP_ALIGN.CENTER, space_after=Pt(6))

tf3 = tb(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.6))
para(tf3, "EDA Course Project  •  April 2026", sz=16, align=PP_ALIGN.CENTER, color=GRAY, first=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Problem Statement")
tf = content(s)
bullet(tf, 'In ASIC technology mapping, the ABC tool generates structural "cuts" in an And-Inverter Graph (AIG) and selects the best ones to minimize area and delay.', sz=19)
bullet(tf, "ABC uses greedy heuristics like Area Flow and MFFC to evaluate cuts, but these are local — they can miss globally optimal solutions.", sz=19)
bullet(tf, "Increasing the cut limit (e.g., from C=8 to C=16) improves Quality of Results (QoR), but nearly doubles the CPU runtime.", bold=True, sz=19)
bullet(tf, "Goal: Can a machine learning model intelligently select superior cuts from the default pool (C=8), achieving the QoR of a larger pool without the runtime penalty?", bold=True, sz=20)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation: Runtime vs Cut Count
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Motivation: Mapping Time vs. Cut Capacity")
tf = content(s)
bullet(tf, "We measured baseline CPU mapping time across 8 benchmarks while varying the max-cuts parameter (-C).", sz=19)
bullet(tf, "As the cut limit increases, the runtime grows significantly while QoR gains diminish.", sz=19)
bullet(tf, "This demonstrates the need for smarter cut selection rather than brute-force enumeration.", sz=19)

data = [
    ["Max Cuts (-C)", "Total CPU Time (s)", "Normalized Time"],
    ["2",  "1.225", "0.37×"],
    ["4",  "1.855", "0.55×"],
    ["6",  "2.489", "0.74×"],
    ["8 (default)",  "3.352", "1.00× (baseline)"],
    ["12", "4.666", "1.39×"],
    ["16", "6.066", "1.81×"],
]
add_table(s, data, Inches(3.2), Inches(3.6), Inches(7), 0.42,
          [2.0, 2.5, 2.5])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Approach Overview
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Our Approach: Overview")
tf = content(s)
bullet(tf, "We trained a lightweight Multi-Layer Perceptron (MLP) to score cuts based on 9 hand-engineered structural features.", sz=19)
bullet(tf, "The trained model was converted into pure C code and integrated directly into ABC's mapping loop — no Python or file I/O at runtime.", sz=19)
bullet(tf, 'Instead of filtering cuts, we use an "Area Blending" strategy: the ML score gently adjusts ABC\'s internal area estimate during its area-recovery rounds.', sz=19)
bullet(tf, "A conservative blending factor (α = 0.02) ensures the ML acts as an intelligent tie-breaker without disrupting ABC's convergence guarantees.", sz=19)
bullet(tf, "This approach allows the ML signal to organically propagate through the mapping graph.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Comparison with SLAP
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Comparison with SLAP Framework")
tf = content(s)
bullet(tf, "SLAP (Supervised Learning for ASIC Performance) is a recent ML-for-mapping framework. Our approach diverges in two key ways:", sz=18)
bullet(tf, "1. Model Architecture", bold=True, sz=19)
bullet(tf, "SLAP uses complex Graph Neural Networks (GNNs) and PairNet structures to capture global circuit topology.", level=1, sz=17)
bullet(tf, "We use a simple 3-layer MLP and compensate with carefully hand-crafted structural features (e.g., mffc_size, slack_ratio) computed natively by ABC.", level=1, sz=17)
bullet(tf, "2. Integration Strategy", bold=True, sz=19)
bullet(tf, "SLAP pre-filters cuts — scoring thousands and discarding weak ones before the mapper runs.", level=1, sz=17)
bullet(tf, "We blend the ML score inline during ABC's dynamic programming phase, adjusting the area estimate in real time.", level=1, sz=17)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — ML Model Architecture
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "ML Model: Architecture")
tf = content(s)
bullet(tf, "The model is a 3-layer feed-forward neural network, implemented in PyTorch and later exported to C.", sz=19)
bullet(tf, "Input Layer:     9 neurons — one per structural feature", level=1, sz=18)
bullet(tf, "Hidden Layer 1:  64 neurons with ReLU activation and 20% Dropout", level=1, sz=18)
bullet(tf, "Hidden Layer 2:  32 neurons with ReLU activation and 20% Dropout", level=1, sz=18)
bullet(tf, "Output Layer:    1 neuron producing a raw scalar quality score", level=1, sz=18)
bullet(tf, "The model was trained on Google Colab using a T4 GPU.", sz=19)
bullet(tf, "After training, the PyTorch weights were exported into a static C header file (model_weights.h) for native execution inside ABC.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Feature Engineering
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Feature Engineering")
tf = content(s)
bullet(tf, "A major challenge was avoiding label leakage. If the model is trained on ABC's own heuristic values (Area Flow, Delay), it simply learns to copy ABC and provides zero improvement.", sz=18)
bullet(tf, "We selected 9 structural/topological features that ABC does not directly optimize:", bold=True, sz=18)
bullet(tf, "n_leaves, node_level, node_fanout — basic cut and node topology", level=1, sz=17)
bullet(tf, "is_critical, slack_ratio — timing pressure indicators", level=1, sz=17)
bullet(tf, "fanout_adj_area, area_per_leaf — structural interaction ratios", level=1, sz=17)
bullet(tf, "mffc_size, mffc_per_leaf — Maximum Fanout-Free Cone efficiency", level=1, sz=17)
bullet(tf, "cut_delay, area_flow, and required_time were used exclusively for constructing the training label and were never passed as input features.", bold=True, sz=18)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Training Strategy
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Label Formulation & Training Strategy")
tf = content(s)
bullet(tf, "Initial Attempt: Labels based on MFFC size alone. Result — exactly 0% QoR improvement because the model simply cloned ABC's existing heuristic.", sz=18)
bullet(tf, "Final Label: We formulated labels based on the Area-Delay Product (ADP):", bold=True, sz=19)
bullet(tf, "Quality Score = −(Area Flow × Cut Delay)", sz=20, bold=True, level=1)
bullet(tf, "A higher quality score means a cut achieves lower area and delay simultaneously.", level=1, sz=18)
bullet(tf, "We used a Pairwise MarginRankingLoss (RankNet approach). The model is shown pairs of cuts for the same node and is penalized if it incorrectly ranks the worse-ADP cut higher.", sz=18)
bullet(tf, "Features are standardized using a fitted StandardScaler before training.", sz=18)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Area Blending
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Integration: Area Blending Architecture")
tf = content(s)
bullet(tf, "Unlike SLAP's pre-filter approach, we inject the ML score directly into ABC's area-recovery mapping rounds (Mode 2).", sz=19)
bullet(tf, "The blending formula modifies ABC's internal area estimate:", sz=19)
para(tf, "Area_adjusted = Area_ABC  ×  (1 − α · tanh(ML_Score))", sz=24, bold=True, align=PP_ALIGN.CENTER, space_after=Pt(14))
bullet(tf, "We used α = 0.02 — a conservative value that ensures the ML acts as a gentle tie-breaker for structurally similar cuts.", sz=19)
bullet(tf, "The ML signal propagates organically through ABC's dynamic programming graph rather than overriding decisions.", sz=19)
bullet(tf, "The C inference engine (ifML.c) executes the MLP forward pass using standard floating-point arithmetic with zero external dependencies.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Pipeline
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "End-to-End Implementation Pipeline")
tf = content(s)
bullet(tf, "Step 00–02: Set up dependencies, instrument ABC to dump structural features, and analyze the profiling data.", sz=18)
bullet(tf, "Step 03–04: Patch the ABC source code and generate training data across all benchmarks.", sz=18)
bullet(tf, "Step 05: Clean data, engineer 9 structural features, and train the MLP on Google Colab (T4 GPU) using RankNet pairwise loss.", sz=18)
bullet(tf, "Step 06: Export the trained PyTorch weights into a static C header file (model_weights.h).", sz=18)
bullet(tf, "Step 07: Inject the ML inference engine and weights into the ABC source tree and recompile.", sz=18)
bullet(tf, "Step 08: Run the final QoR comparison between baseline ABC and ML-guided ABC on all benchmarks.", sz=18)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Challenges
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Challenges Faced & Debugging")
tf = content(s)
bullet(tf, 'Challenge 1: The "Trivial Cut" Crashing Bug', bold=True, sz=19)
bullet(tf, "The ML model was initially allowed to force ABC to choose its highest-scoring cut after mapping.", level=1, sz=17)
bullet(tf, "The model assigned extremely high scores to trivial cuts (single-input identity cuts with near-zero delay and area).", level=1, sz=17)
bullet(tf, "This fractured the mapping graph, causing massive LUT bloat and segmentation faults.", level=1, sz=17)
bullet(tf, "Resolution: Filter out trivial cuts (n_leaves < 2) before ML scoring.", level=1, sz=17, bold=True)
bullet(tf, 'Challenge 2: The "Zero Improvement" Phase', bold=True, sz=19)
bullet(tf, "Even after fixing crashes, the model yielded exactly 0% difference from baseline ABC.", level=1, sz=17)
bullet(tf, "Root causes: Label leakage (training on MFFC) and overriding cuts after convergence broke timing.", level=1, sz=17)
bullet(tf, "Resolution: Switched to ADP-based labels and inline Area Blending architecture.", level=1, sz=17, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Results (Arithmetic)
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Results: QoR Comparison (Arithmetic Benchmarks)")
tf = content(s)
bullet(tf, "Comparison of baseline ABC (C=8) vs ML-Guided ABC on complex arithmetic circuits.", sz=17)

results_arith = [
    ["Benchmark",   "ABC LUTs", "ML LUTs", "LUT Imp %", "ABC ADP",     "ML ADP",      "ADP Imp %"],
    ["adder",       "254",      "254",     "0.00",      "12,954",      "12,954",      "0.00"],
    ["bar",         "1,456",    "1,456",   "0.00",      "43,680",      "43,680",      "0.00"],
    ["div",         "22,030",   "19,942",  "+9.48",     "19,033,920",  "17,229,888",  "+9.48"],
    ["hyp",         "44,508",   "41,969",  "+5.70",     "186,666,552", "176,017,986", "+5.70"],
    ["log2",        "8,010",    "7,778",   "+2.90",     "616,770",     "598,906",     "+2.90"],
    ["max",         "842",      "831",     "+1.31",     "47,152",      "46,536",      "+1.31"],
    ["multiplier",  "5,929",    "5,683",   "+4.15",     "314,237",     "301,199",     "+4.15"],
    ["sin",         "1,464",    "1,421",   "+2.94",     "61,488",      "59,682",      "+2.94"],
    ["sqrt",        "5,724",    "5,180",   "+9.50",     "5,912,892",   "5,350,940",   "+9.50"],
    ["square",      "3,997",    "3,502",   "+12.38",    "199,850",     "175,100",     "+12.38"],
]
add_table(s, results_arith, Inches(0.5), Inches(2.1), Inches(12.3), 0.38,
          [1.8, 1.4, 1.4, 1.4, 2.1, 2.1, 1.4], header=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Results (Random Control)
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Results: QoR Comparison (Random Control Benchmarks)")
tf = content(s)
bullet(tf, "Comparison of baseline ABC (C=8) vs ML-Guided ABC on random control circuits.", sz=17)

results_ctrl = [
    ["Benchmark",   "ABC LUTs", "ML LUTs", "LUT Imp %", "ABC ADP",     "ML ADP",      "ADP Imp %"],
    ["cavlc",       "359",      "359",     "0.00",      "2,513",       "2,513",       "0.00"],
    ["ctrl",        "57",       "57",      "0.00",      "456",         "456",         "0.00"],
    ["dec",         "304",      "304",     "0.00",      "608",         "608",         "0.00"],
    ["i2c",         "365",      "362",     "+0.82",     "1,460",       "1,448",       "+0.82"],
    ["int2float",   "123",      "123",     "0.00",      "861",         "861",         "0.00"],
    ["mem_ctrl",    "12,104",   "11,914",  "+1.57",     "302,600",     "297,850",     "+1.57"],
    ["priority",    "219",      "218",     "+0.46",     "6,789",       "6,758",       "+0.46"],
    ["router",      "105",      "105",     "0.00",      "1,575",       "1,575",       "0.00"],
    ["voter",       "2,826",    "2,785",   "+1.45",     "48,042",      "47,345",      "+1.45"],
]
add_table(s, results_ctrl, Inches(0.5), Inches(2.1), Inches(12.3), 0.38,
          [1.8, 1.4, 1.4, 1.4, 2.1, 2.1, 1.4], header=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Results Summary
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Results: Summary & Analysis")
tf = content(s)
bullet(tf, "The ML-guided mapper was evaluated across 19 AIGER benchmarks.", sz=19)
bullet(tf, "Average ADP Reduction (arithmetic mean): 2.63%", bold=True, sz=21)
bullet(tf, "Zero instances of QoR degradation across all benchmarks.", bold=True, sz=19)
bullet(tf, "Best-performing benchmarks:", bold=True, sz=19)
bullet(tf, "square (multiplier):  12.38% ADP improvement — largest gain observed", level=1, sz=18)
bullet(tf, "sqrt:                        9.50% ADP improvement", level=1, sz=18)
bullet(tf, "div (divider):            9.48% ADP improvement", level=1, sz=18)
bullet(tf, "hyp (hyperbolic):       5.70% ADP improvement", level=1, sz=18)
bullet(tf, "The model shows strongest gains on complex arithmetic circuits where ABC's local heuristics are most likely to miss globally optimal cut selections.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Codebase Structure
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Codebase Structure")
tf = content(s)
bullet(tf, "abc/ — The modified ABC synthesis tool with our ML integration hooks.", sz=18)
bullet(tf, "benchmarks/ — AIGER benchmark circuits (arithmetic and random control suites).", sz=18)
bullet(tf, "ml_cut_project/ — Main ML pipeline directory containing:", sz=18)
bullet(tf, "ml/ — Trained model weights (cut_model_mlp.pt) and feature scaler (scaler.pkl).", level=1, sz=17)
bullet(tf, "abc_patch/ — Source patches: ifML.c (inference engine) and model_weights.h.", level=1, sz=17)
bullet(tf, "data/ and results/ — Raw CSV dumps and QoR comparison logs.", level=1, sz=17)
bullet(tf, "Core ABC Modifications:", bold=True, sz=18)
bullet(tf, "ifML.c — Contains the If_MLBlendCutArea function that runs the MLP forward pass.", level=1, sz=17)
bullet(tf, "ifMap.c — Modified to call the ML blending function during area-recovery rounds.", level=1, sz=17)
bullet(tf, "if.h — Updated header with forward declarations for the ML integration.", level=1, sz=17)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Key Takeaways")
tf = content(s)
bullet(tf, "A lightweight MLP can outperform ABC's highly optimized, hand-tuned EDA heuristics when guided by the right features and labels.", sz=19)
bullet(tf, "Avoiding label leakage is critical — training on structural features forces the model to learn genuinely new signals rather than copying ABC.", sz=19)
bullet(tf, "Area Blending is superior to hard filtering because it preserves the optimizer's mathematical convergence guarantees.", sz=19)
bullet(tf, "Native C integration (model_weights.h + ifML.c) eliminates all I/O overhead, making ML inference practical inside tight mapping loops.", sz=19)
bullet(tf, "The results demonstrate that ML can act not just as an external filter, but as an embedded mathematical guide for traditional EDA algorithms.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 17 — Conclusion
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); heading(s, "Conclusion")
tf = content(s)
bullet(tf, "This project successfully implemented an ML-guided ASIC technology mapper by deeply integrating a structural-feature-based MLP directly into ABC's C codebase.", sz=19)
bullet(tf, "We overcame significant integration hurdles — trivial cut crashes, label leakage, and broken timing graphs — through systematic debugging.", sz=19)
bullet(tf, "The Area Blending architecture proved to be the correct design choice, allowing ML to gently steer ABC's decisions.", sz=19)
bullet(tf, "The final model achieved up to 12.38% ADP reduction on complex arithmetic circuits, with an average improvement of 2.63% and zero regressions.", bold=True, sz=19)
bullet(tf, "These results highlight the practical potential of machine learning for improving EDA tool quality without sacrificing runtime or stability.", sz=19)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 18 — Thank You
# ═══════════════════════════════════════════════════════════════════
s = add_slide(); deco_bars(s)
tf = tb(s, Inches(1), Inches(2.2), Inches(11.3), Inches(3))
para(tf, "Thank You", sz=48, bold=True, align=PP_ALIGN.CENTER, space_after=Pt(24), first=True)
para(tf, "Questions & Discussion", sz=28, align=PP_ALIGN.CENTER, color=GRAY)

# ═══════════════════════════════════════════════════════════════════
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ML_Cut_Selection_Presentation.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
